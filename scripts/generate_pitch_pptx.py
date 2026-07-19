from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "public" / "SurgiFind-Pitch-Deck.pptx"
TEAM_DIR = ROOT / "public" / "team"

NAVY = RGBColor(10, 22, 38)
SLATE = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
MIST = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(5, 171, 165)
TEAL_DARK = RGBColor(13, 116, 111)
TEAL_LIGHT = RGBColor(236, 253, 245)
CYAN_TINT = RGBColor(236, 254, 255)
GOLD = RGBColor(245, 158, 11)
ROSE = RGBColor(244, 63, 94)
BORDER = RGBColor(203, 213, 225)

TEAM = [
    ("Aditi Prasad", "SDE & Product Lead", TEAM_DIR / "aditi.png"),
    ("Mayank Prasad", "SDE", TEAM_DIR / "mayank.jpg"),
    ("Deepnarayan Lohra", "SDE", TEAM_DIR / "deepnarayan.png"),
]


def add_rect(slide, left, top, width, height, fill, line=None, radius=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=SLATE, bold=False,
             align=PP_ALIGN.LEFT, font_name="Aptos", italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    run.font.italic = italic
    return box


def add_bullets(slide, left, top, width, height, items, size=20, color=SLATE, bullet_color=TEAL):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        if p.runs:
            p.runs[0].font.color.rgb = color
    return box


def add_chip(slide, left, top, width, text, fill=TEAL, text_color=WHITE):
    add_rect(slide, left, top, width, 0.42, fill, line=fill)
    add_text(slide, left, top + 0.03, width, 0.28, text, size=11, color=text_color, bold=True, align=PP_ALIGN.CENTER)


def add_top_label(slide, label, dark=False):
    fill = CYAN_TINT if not dark else RGBColor(18, 45, 66)
    line = RGBColor(165, 243, 252) if not dark else RGBColor(45, 212, 191)
    text_color = TEAL_DARK if not dark else RGBColor(153, 246, 228)
    add_rect(slide, 0.55, 0.42, 2.05, 0.5, fill, line=line)
    add_text(slide, 0.72, 0.54, 1.7, 0.2, label, size=11, color=text_color, bold=True)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    add_rect(slide, 0.0, 0.0, 13.33, 7.5, NAVY)
    add_rect(slide, 8.7, -0.6, 5.6, 4.0, RGBColor(17, 94, 89), line=None, radius=MSO_SHAPE.OVAL)
    add_rect(slide, 10.1, 4.7, 3.3, 2.5, RGBColor(8, 47, 73), line=None, radius=MSO_SHAPE.OVAL)
    add_top_label(slide, "HACKATHON PITCH", dark=True)
    add_text(slide, 0.7, 1.25, 6.3, 0.8, "SurgiFind", size=28, color=WHITE, bold=True)
    add_text(slide, 0.7, 2.0, 7.5, 1.5, "From surgical uncertainty to hospital booking in one guided flow.", size=28, color=WHITE, bold=True)
    add_text(slide, 0.72, 3.45, 5.5, 1.0, "Search. Compare price. Check insurance. Book with confidence.", size=17, color=RGBColor(203, 213, 225))

    card_data = [
        ("Search", "Surgery, city, budget, rating"),
        ("Trust", "Price and insurance clarity"),
        ("Action", "Authenticated booking and history"),
    ]
    left = 0.72
    for title, body in card_data:
        add_rect(slide, left, 5.35, 2.2, 1.15, RGBColor(255, 255, 255), line=RGBColor(45, 212, 191))
        add_text(slide, left + 0.15, 5.58, 1.8, 0.26, title, size=15, color=SLATE, bold=True)
        add_text(slide, left + 0.15, 5.93, 1.85, 0.38, body, size=10, color=MUTED)
        left += 2.45

    add_rect(slide, 7.55, 1.35, 4.95, 4.95, RGBColor(248, 250, 252), line=RGBColor(153, 246, 228))
    add_text(slide, 8.0, 1.8, 4.05, 0.4, "Why this matters", size=14, color=TEAL_DARK, bold=True)
    add_text(slide, 8.0, 2.25, 3.9, 2.15, "A patient trying to book surgery today jumps across hospital sites, price ambiguity, insurance confusion, and manual calls. SurgiFind compresses that chaos into one decision flow.", size=21, color=SLATE, bold=True)
    add_text(slide, 8.0, 4.75, 3.7, 0.65, "Built as a working full-stack prototype with AI-assisted discovery, secure booking, and a production-ready path forward.", size=12, color=MUTED)
    add_chip(slide, 10.85, 6.42, 1.05, "01")
    return slide


def problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_top_label(slide, "THE PAIN")
    add_text(slide, 0.7, 1.0, 6.5, 0.9, "Patients do not need more hospital tabs. They need a decision path.", size=28, color=SLATE, bold=True)
    add_rect(slide, 0.75, 1.95, 5.0, 4.65, CYAN_TINT, line=RGBColor(165, 243, 252))
    add_text(slide, 1.05, 2.3, 4.2, 0.35, "Real-world moment", size=13, color=TEAL_DARK, bold=True)
    add_text(slide, 1.05, 2.8, 4.1, 1.5, "A family knows the surgery type, but not which hospital is trustworthy, affordable, in-network, or even bookable today.", size=24, color=SLATE, bold=True)
    add_text(slide, 1.05, 4.65, 4.0, 1.15, "What follows is fragmented research, hidden pricing, insurance guesswork, and repeated manual calls before any real progress happens.", size=15, color=MUTED)

    pain_cards = [
        ("Opaque pricing", "Patients cannot compare likely surgery costs with confidence.", GOLD),
        ("Insurance friction", "Network and coverage checks take too long and break decision flow.", ROSE),
        ("No clean booking path", "Discovery and booking are disconnected, so conversion stalls.", TEAL),
    ]
    top = 2.0
    for title, body, accent in pain_cards:
        add_rect(slide, 6.15, top, 6.0, 1.28, WHITE, line=BORDER)
        add_rect(slide, 6.35, top + 0.2, 0.18, 0.88, accent, line=accent, radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(slide, 6.7, top + 0.2, 2.2, 0.28, title, size=17, color=SLATE, bold=True)
        add_text(slide, 6.7, top + 0.55, 4.9, 0.48, body, size=12, color=MUTED)
        top += 1.52

    add_chip(slide, 11.8, 0.55, 0.95, "02")
    return slide


def solution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    add_top_label(slide, "THE PRODUCT")
    add_text(slide, 0.7, 1.0, 6.0, 0.9, "SurgiFind turns surgery discovery into one connected workflow.", size=28, color=SLATE, bold=True)
    add_text(slide, 0.7, 1.72, 5.8, 0.65, "Instead of asking patients to stitch together research, cost checks, insurance lookup, and booking, we orchestrate it in one place.", size=15, color=MUTED)

    add_rect(slide, 0.72, 2.55, 5.2, 3.85, WHITE, line=BORDER)
    add_text(slide, 1.02, 2.9, 4.5, 0.35, "The wedge", size=13, color=TEAL_DARK, bold=True)
    add_text(slide, 1.02, 3.3, 4.35, 1.0, "Search + transparency + booking", size=26, color=SLATE, bold=True)
    add_bullets(slide, 1.02, 4.35, 4.2, 1.7, [
        "Hospital and surgery search with filters",
        "Price and insurance context before action",
        "Chat guidance for natural-language discovery",
        "Authenticated booking and booking history",
    ], size=15)

    flow = [
        ("Search", "Surgery, city, budget, rating"),
        ("Compare", "Hospital options and price ranges"),
        ("Ask", "Chat assistant for guided discovery"),
        ("Book", "Reserve a consultation or surgery slot"),
    ]
    left = 6.35
    for index, (title, body) in enumerate(flow):
        add_rect(slide, left, 2.85, 1.45, 2.55, WHITE, line=RGBColor(153, 246, 228))
        add_chip(slide, left + 0.34, 3.08, 0.78, f"0{index + 1}")
        add_text(slide, left + 0.18, 3.78, 1.05, 0.35, title, size=18, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, left + 0.12, 4.25, 1.18, 0.75, body, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        if index < len(flow) - 1:
            add_text(slide, left + 1.48, 4.02, 0.35, 0.25, "→", size=24, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
        left += 1.62

    add_chip(slide, 11.8, 0.55, 0.95, "03")
    return slide


def demo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_top_label(slide, "LIVE DEMO ARC", dark=True)
    add_text(slide, 0.7, 1.0, 6.5, 0.9, "What judges will see in under two minutes", size=28, color=WHITE, bold=True)
    add_text(slide, 0.7, 1.7, 6.2, 0.6, "A clean flow beats feature overload. The story is discovery to action.", size=15, color=RGBColor(191, 219, 254))

    demo_steps = [
        ("1", "Search", "Find surgery options by city, budget, and rating."),
        ("2", "Hospital", "Open detail page and show pricing + slot context."),
        ("3", "Chat", "Ask in plain language and get guided matches."),
        ("4", "Book", "Confirm a slot, then show booking history and cancellation."),
    ]
    left = 0.8
    for number, title, body in demo_steps:
        add_rect(slide, left, 2.65, 2.9, 2.85, RGBColor(15, 34, 54), line=RGBColor(45, 212, 191))
        add_chip(slide, left + 0.18, 2.9, 0.6, number, fill=WHITE, text_color=TEAL_DARK)
        add_text(slide, left + 0.18, 3.45, 2.2, 0.3, title, size=20, color=WHITE, bold=True)
        add_text(slide, left + 0.18, 3.95, 2.35, 0.95, body, size=13, color=RGBColor(203, 213, 225))
        if left < 9.3:
            add_text(slide, left + 2.98, 3.83, 0.3, 0.25, "→", size=24, color=RGBColor(94, 234, 212), bold=True, align=PP_ALIGN.CENTER)
        left += 3.15

    add_rect(slide, 0.82, 5.9, 11.65, 0.9, RGBColor(9, 29, 46), line=RGBColor(22, 78, 99))
    add_text(slide, 1.08, 6.17, 10.9, 0.25, "Key message: we are not showing isolated pages; we are showing a patient journey that reaches a real booking outcome.", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_chip(slide, 11.8, 0.55, 0.95, "04", fill=RGBColor(94, 234, 212), text_color=SLATE)
    return slide


def advantage_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_top_label(slide, "WHY THIS CAN WIN")
    add_text(slide, 0.7, 1.0, 6.2, 0.9, "Most tools stop at discovery. SurgiFind continues to decision and action.", size=28, color=SLATE, bold=True)

    columns = [
        ("Discovery", "Search hospitals by surgery, city, budget, rating, and type.", CYAN_TINT),
        ("Decision", "Give price clarity, insurance relevance, and conversational support.", TEAL_LIGHT),
        ("Action", "Convert intent into secure booking, history, and cancellation.", RGBColor(255, 247, 237)),
    ]
    left = 0.82
    for title, body, fill in columns:
        add_rect(slide, left, 2.15, 3.85, 3.25, fill, line=BORDER)
        add_text(slide, left + 0.22, 2.47, 2.8, 0.35, title, size=21, color=SLATE, bold=True)
        add_text(slide, left + 0.22, 3.0, 3.1, 1.0, body, size=14, color=MUTED)
        left += 4.12

    add_rect(slide, 0.82, 5.75, 12.0, 0.9, SLATE, line=SLATE)
    add_text(slide, 1.15, 6.03, 11.25, 0.23, "This is a workflow product, not just an information product. That is the core differentiator.", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_chip(slide, 11.8, 0.55, 0.95, "05")
    return slide


def stakeholders_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    add_top_label(slide, "VALUE CREATION")
    add_text(slide, 0.7, 1.0, 5.5, 0.9, "Why this is bigger than a demo", size=28, color=SLATE, bold=True)
    add_text(slide, 0.7, 1.72, 5.4, 0.65, "SurgiFind creates value for the three parties that matter in planned care.", size=15, color=MUTED)

    cards = [
        ("Patients", ["Less confusion before surgery", "Faster shortlist creation", "Clarity on cost and booking"]),
        ("Hospitals", ["Qualified high-intent leads", "Better booking conversion", "Lower manual follow-up load"]),
        ("Partners", ["Insurance relevance", "Future referral workflows", "Operational integration potential"]),
    ]
    left = 0.82
    for title, items in cards:
        add_rect(slide, left, 2.6, 3.86, 3.2, WHITE, line=BORDER)
        add_text(slide, left + 0.24, 2.95, 2.5, 0.32, title, size=21, color=SLATE, bold=True)
        add_bullets(slide, left + 0.18, 3.45, 3.2, 1.9, items, size=15)
        left += 4.08

    add_rect(slide, 0.82, 6.0, 12.0, 0.62, CYAN_TINT, line=RGBColor(165, 243, 252))
    add_text(slide, 1.1, 6.18, 11.3, 0.2, "Revenue paths: lead generation, premium listings, booking commissions, and care navigation partnerships.", size=14, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_chip(slide, 11.8, 0.55, 0.95, "06")
    return slide


def scale_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_top_label(slide, "NEXT 90 DAYS", dark=True)
    add_text(slide, 0.7, 1.0, 6.2, 0.9, "We know exactly what comes after the prototype", size=28, color=WHITE, bold=True)
    add_text(slide, 0.7, 1.7, 5.8, 0.6, "The path forward is payments, production cloud operations, and stronger platform controls.", size=15, color=RGBColor(191, 219, 254))

    milestones = [
        ("Payments", "Deposits, confirmations, refunds, receipts"),
        ("Cloud", "Vercel deployment, CI/CD, managed secrets"),
        ("Security", "Rate limiting, WAF, audit logging"),
        ("Scale", "Cache, queue, autoscaling, observability"),
    ]
    left = 0.82
    for title, body in milestones:
        add_rect(slide, left, 2.45, 2.82, 1.75, RGBColor(15, 34, 54), line=RGBColor(45, 212, 191))
        add_text(slide, left + 0.18, 2.78, 1.8, 0.3, title, size=18, color=WHITE, bold=True)
        add_text(slide, left + 0.18, 3.25, 2.2, 0.75, body, size=12, color=RGBColor(203, 213, 225))
        left += 3.05

    add_rect(slide, 0.82, 4.7, 12.0, 1.7, WHITE, line=RGBColor(45, 212, 191))
    labels = [
        ("Users", 1.1),
        ("Edge", 3.25),
        ("App", 5.4),
        ("Platform", 7.55),
        ("Data + Integrations", 9.7),
    ]
    for label, left in labels:
        add_rect(slide, left, 5.12, 1.7, 0.8, TEAL_LIGHT if label in {"Edge", "Platform"} else MIST, line=RGBColor(153, 246, 228))
        add_text(slide, left + 0.05, 5.34, 1.6, 0.2, label, size=13, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    for left in [2.86, 5.01, 7.16, 9.31]:
        add_text(slide, left, 5.34, 0.28, 0.18, "→", size=22, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)

    add_chip(slide, 11.8, 0.55, 0.95, "07", fill=RGBColor(94, 234, 212), text_color=SLATE)
    return slide


def team_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_top_label(slide, "TEAM + CLOSE")
    add_text(slide, 0.7, 1.0, 6.0, 0.9, "Built by a team thinking about product, not just code", size=28, color=SLATE, bold=True)
    add_text(slide, 0.7, 1.72, 5.9, 0.6, "We built the workflow end-to-end: search, AI guidance, booking, history, cancellation, deployment, and the scale-up plan.", size=15, color=MUTED)

    left = 0.82
    for name, role, image_path in TEAM:
        add_rect(slide, left, 2.45, 3.85, 3.05, WHITE, line=BORDER)
        if image_path.exists():
            slide.shapes.add_picture(str(image_path), Inches(left + 0.26), Inches(2.72), width=Inches(3.32), height=Inches(1.82))
        add_text(slide, left + 0.26, 4.72, 3.0, 0.28, name, size=16, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, left + 0.26, 5.04, 3.0, 0.22, role, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        left += 4.06

    add_rect(slide, 0.82, 5.95, 12.0, 0.75, SLATE, line=SLATE)
    add_text(slide, 1.1, 6.18, 11.2, 0.18, "SurgiFind makes surgery discovery transparent and patient-first. That is the product, and that is the pitch.", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_chip(slide, 11.8, 0.55, 0.95, "08")
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    problem_slide(prs)
    solution_slide(prs)
    demo_slide(prs)
    advantage_slide(prs)
    stakeholders_slide(prs)
    scale_slide(prs)
    team_slide(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    main()
